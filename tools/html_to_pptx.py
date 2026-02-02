from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from bs4 import BeautifulSoup
import os

INPUT_HTML = "presentation.html"
OUTPUT_PPTX = "presentation.pptx"


def px_in_inches(px, dpi=96):
    return px / dpi


def add_shadow(shape):
    """Add subtle shadow to shape"""
    try:
        shadow = shape.shadow
        shadow.inherit = False
        shadow.angle = 90
        shadow.blur_radius = Pt(6)
        shadow.distance = Pt(2)
        shadow.alpha = 0.15
    except:
        pass


def parse_html(path):
    with open(path, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f, "lxml")

    items = []
    for div in soup.select(".timeline-item"):
        side = "left" if "left" in div.get("class", []) else "right"
        card = div.select_one(".card")
        date = (
            card.select_one(".timeline-date").get_text(strip=True)
            if card and card.select_one(".timeline-date")
            else ""
        )
        title = (
            card.select_one(".timeline-title").get_text(strip=True)
            if card and card.select_one(".timeline-title")
            else ""
        )
        subtitle = (
            card.select_one(".timeline-subtitle").get_text(strip=True)
            if card and card.select_one(".timeline-subtitle")
            else ""
        )
        is_current = bool(
            div.select_one(".current-dot")
            or card
            and "current-role" in card.get("class", [])
        )
        items.append(
            {
                "side": side,
                "date": date,
                "title": title,
                "subtitle": subtitle,
                "current": is_current,
            }
        )

    # Role card
    role = {}
    role_card = soup.select_one(".role-card")
    if role_card:
        # title may be in h3
        h3 = role_card.find("h3")
        role["title"] = h3.get_text(strip=True) if h3 else ""
        # first small p date
        pdate = role_card.find("p")
        role["date"] = pdate.get_text(strip=True) if pdate else ""
        # description: first big paragraph after header
        paras = role_card.find_all("p")
        role["desc"] = paras[1].get_text(strip=True) if len(paras) > 1 else ""
        # responsibilities
        resp = []
        for li in role_card.select("ul li"):
            txt = li.get_text(separator=" ", strip=True)
            resp.append(txt)
        role["items"] = resp
    return items, role


def create_pptx(items, role):
    prs = Presentation()
    prs.slide_width = Inches(px_in_inches(1280))
    prs.slide_height = Inches(px_in_inches(720))

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width
    H = prs.slide_height

    # Header: dark blue (#0f172a) with white text
    header_h = Inches(0.7)
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, W, header_h)
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(15, 23, 42)
    header.line.fill.background()

    tf = header.text_frame
    tf.margin_left = Inches(0.24)
    tf.margin_top = Inches(0.12)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Background"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf.add_paragraph()
    p2.text = "Timeline, Current Role & UAV Gear"
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(191, 219, 254)
    p2.space_before = Pt(2)

    # Timeline column: left 65%
    timeline_w = W * 0.65
    timeline_left = 0
    center_x = timeline_w * 0.5

    # Vertical timeline line
    line_w = Pt(4)
    line_left = center_x - line_w / 2
    line_top = header_h + Inches(0.2)
    line_h = H - header_h - Inches(0.4)

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        int(line_left),
        int(line_top),
        int(line_w),
        int(line_h),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(226, 232, 240)
    line.line.fill.background()

    # Timeline items
    n = len(items)
    top_margin = header_h + Inches(0.3)
    bottom_margin = Inches(0.3)
    usable_h = H - top_margin - bottom_margin
    item_spacing = usable_h / n

    for i, item in enumerate(items):
        y_center = top_margin + (i + 0.5) * item_spacing

        # Dot
        dot_size = Pt(16)
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            int(center_x - dot_size / 2),
            int(y_center - dot_size / 2),
            int(dot_size),
            int(dot_size),
        )
        dot.fill.solid()
        if item["current"]:
            dot.fill.fore_color.rgb = RGBColor(245, 158, 11)
            dot.line.color.rgb = RGBColor(252, 211, 77)
        else:
            dot.fill.fore_color.rgb = RGBColor(37, 99, 235)
            dot.line.color.rgb = RGBColor(191, 219, 254)
        dot.line.width = Pt(3)

        # Card
        card_h = Inches(0.85)
        card_w_pct = 0.42

        if item["side"] == "left":
            card_w = timeline_w * card_w_pct
            card_left = center_x - card_w - Inches(0.35)
        else:
            card_w = timeline_w * card_w_pct
            card_left = center_x + Inches(0.35)

        card_top = y_center - card_h / 2

        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            int(card_left),
            int(card_top),
            int(card_w),
            int(card_h),
        )
        card.fill.solid()
        if item["current"]:
            card.fill.fore_color.rgb = RGBColor(239, 246, 255)
        else:
            card.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Border
        if item["side"] == "left":
            card.line.color.rgb = (
                RGBColor(37, 99, 235) if not item["current"] else RGBColor(245, 158, 11)
            )
            card.line.width = Pt(4)
        else:
            card.line.color.rgb = (
                RGBColor(37, 99, 235) if not item["current"] else RGBColor(245, 158, 11)
            )
            card.line.width = Pt(4)

        add_shadow(card)

        # Card text
        tf = card.text_frame
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(12)
        tf.margin_top = Pt(8)
        tf.margin_bottom = Pt(8)
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = item["date"]
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = (
            RGBColor(37, 99, 235) if not item["current"] else RGBColor(245, 158, 11)
        )
        p.space_after = Pt(2)

        p2 = tf.add_paragraph()
        p2.text = item["title"]
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(31, 41, 55)
        p2.space_after = Pt(1)

        p3 = tf.add_paragraph()
        p3.text = item["subtitle"]
        p3.font.size = Pt(9)
        p3.font.italic = True
        p3.font.color.rgb = RGBColor(75, 85, 99)

    # Right column: Role card and PhD motivation
    right_left = timeline_w + Inches(0.1)
    right_w = W - right_left - Inches(0.15)

    # Background for right column
    right_bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        int(right_left - Inches(0.05)),
        int(header_h),
        int(right_w + Inches(0.2)),
        int(H - header_h),
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(250, 251, 252)
    right_bg.line.fill.background()
    right_bg.shadow.inherit = False

    # Role card
    role_top = header_h + Inches(0.2)
    role_h = Inches(2.8)

    role_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        int(right_left),
        int(role_top),
        int(right_w),
        int(role_h),
    )
    role_card.fill.solid()
    role_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    role_card.line.color.rgb = RGBColor(245, 158, 11)
    role_card.line.width = Pt(3)
    add_shadow(role_card)

    tf = role_card.text_frame
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(12)
    tf.margin_bottom = Pt(12)
    tf.word_wrap = True

    # Title right-aligned
    p = tf.paragraphs[0]
    p.text = role.get("title", "Research Fellow — ISTC CNR")
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)
    p.alignment = PP_ALIGN.RIGHT

    p2 = tf.add_paragraph()
    p2.text = role.get("date", "Jan 2024 – Present")
    p2.font.size = Pt(9)
    p2.font.color.rgb = RGBColor(100, 116, 139)
    p2.alignment = PP_ALIGN.RIGHT
    p2.space_after = Pt(8)

    p3 = tf.add_paragraph()
    p3.text = role.get("desc", "")
    p3.font.size = Pt(8)
    p3.font.color.rgb = RGBColor(71, 85, 105)
    p3.space_after = Pt(6)

    # Responsibilities header
    p4 = tf.add_paragraph()
    p4.text = "RESPONSIBILITIES"
    p4.font.size = Pt(8)
    p4.font.bold = True
    p4.font.color.rgb = RGBColor(51, 65, 85)
    p4.space_after = Pt(4)

    items_text = role.get("items", [])
    if items_text:
        for idx, it in enumerate(items_text[:4]):  # limit
            clean = it.replace("•", "").strip()
            if "DJI" in clean or "UviFy" in clean:
                continue
            p5 = tf.add_paragraph()
            p5.text = "• " + clean
            p5.font.size = Pt(7.5)
            p5.font.color.rgb = RGBColor(71, 85, 105)
            p5.level = 0

    # PhD Motivation section below role card
    phd_top = role_top + role_h + Inches(0.15)
    phd_h = H - phd_top - Inches(0.15)

    phd_card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        int(right_left),
        int(phd_top),
        int(right_w),
        int(phd_h),
    )
    phd_card.fill.solid()
    phd_card.fill.fore_color.rgb = RGBColor(254, 249, 239)
    phd_card.line.color.rgb = RGBColor(245, 158, 11)
    phd_card.line.width = Pt(2)
    add_shadow(phd_card)

    tf = phd_card.text_frame
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)
    tf.margin_bottom = Pt(10)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "PhD Motivation"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(120, 53, 15)
    p.space_after = Pt(6)

    motivation = [
        (
            "Challenge",
            "Real-world autonomy faces energy, perception, and uncertainty constraints",
        ),
        ("Gap", "Classical pipelines reach fundamental limits"),
        (
            "Inspiration",
            "Biological systems (bees) achieve efficient, adaptive navigation",
        ),
        ("Direction", "Neuro-inspired and neuromorphic computation"),
        ("Goal", "Principled foundations for robust, efficient autonomy"),
    ]

    for label, text in motivation:
        p_label = tf.add_paragraph()
        p_label.text = f"{label}: {text}"
        p_label.font.size = Pt(7.5)
        p_label.font.color.rgb = RGBColor(120, 53, 15)
        p_label.space_after = Pt(3)
        # Bold the label
        run = p_label.runs[0]
        run.text = label + ": "
        run.font.bold = True
        # Regular text
        run2 = p_label.add_run()
        run2.text = text
        run2.font.bold = False

    prs.save(OUTPUT_PPTX)
    print("Saved", OUTPUT_PPTX)


if __name__ == "__main__":
    if not os.path.exists(INPUT_HTML):
        print("Input file not found:", INPUT_HTML)
    else:
        items, role = parse_html(INPUT_HTML)
        create_pptx(items, role)
